import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

def create_excel_report(income_df, cash_flows, valuation, scenario_df, sensitivity_df, full_model):
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        # Write 3 statements
        pd.DataFrame(full_model["income_statement"]).to_excel(writer, index=False, sheet_name="Income Statement")
        pd.DataFrame(full_model["cash_flow"]).to_excel(writer, index=False, sheet_name="Cash Flow")
        pd.DataFrame(full_model["balance_sheet"]).to_excel(writer, index=False, sheet_name="Balance Sheet")

        # Valuation and Scenarios
        pd.DataFrame({"Enterprise Value": [valuation]}).to_excel(writer, index=False, sheet_name="Valuation")
        scenario_df.to_excel(writer, index=False, sheet_name="Scenarios")
        sensitivity_df.to_excel(writer, sheet_name="Sensitivity")
    return buffer.getvalue()

def create_powerpoint_report(income_df, valuation, scenario_df):
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Financial Model Report"
    slide.placeholders[1].text = f"Enterprise Value: ${valuation:,.0f}"

    # P&L Slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Income Statement"
    rows, cols = income_df.shape
    table = slide2.shapes.add_table(rows+1, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
    for j, col_name in enumerate(income_df.columns):
        table.cell(0, j).text = col_name
    for i in range(rows):
        for j in range(cols):
            table.cell(i+1, j).text = str(round(income_df.iloc[i, j], 2))

    # Scenarios Slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Scenario Analysis"
    rows, cols = scenario_df.shape
    table2 = slide3.shapes.add_table(rows+1, cols, Inches(1), Inches(2), Inches(8), Inches(3)).table
    for j, col_name in enumerate(scenario_df.columns):
        table2.cell(0, j).text = col_name
    for i in range(rows):
        for j in range(cols):
            table2.cell(i+1, j).text = str(scenario_df.iloc[i, j])

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
